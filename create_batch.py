#!/usr/bin/env python3
"""Create JSONL batch requests with support for multiple document formats.

Supported formats: PDF, DOCX, PPTX, ODP, TXT, MD
"""

import json
import os
import argparse
from pathlib import Path
from pypdf import PdfReader
import pdfplumber
from docx import Document
from pptx import Presentation
from odf.opendocument import load as load_odf
from odf.text import P
from odf.draw import Frame
import glob
from dotenv import load_dotenv
from datetime import datetime

# Parse command line arguments
parser = argparse.ArgumentParser(
    description='Create JSONL batch requests from documents',
    formatter_class=argparse.RawDescriptionHelpFormatter,
    epilog='''
Examples:
  # Process all files in default directory (data/papers/)
  python create_batch.py

  # Process specific files
  python create_batch.py --files paper1.pdf paper2.txt report.docx

  # Process all files in a custom directory
  python create_batch.py --input-dir /path/to/documents/
'''
)
parser.add_argument(
    '--files',
    nargs='+',
    metavar='FILE',
    help='Specific file paths to process'
)
parser.add_argument(
    '--input-dir',
    metavar='DIR',
    help='Directory to scan for documents (default: data/papers/)'
)

args = parser.parse_args()

# Load environment variables
load_dotenv()

# Read summarization prompt and substitute word count
with open('summarisation_prompt.txt', 'r') as f:
    prompt_template = f.read()

# Substitute word count from environment variable (default to 2000)
word_count = os.getenv('SUMMARY_WORD_COUNT', '2000')
prompt_template = prompt_template.replace('{WORD_COUNT}', word_count)

# Print environment variables being used
print("Environment Variables:")
print(f"  SUMMARY_WORD_COUNT: {word_count}")
print(f"  CHAT_COMPLETIONS_ENDPOINT: {os.getenv('CHAT_COMPLETIONS_ENDPOINT', '/v1/chat/completions')}")
print(f"  DOUBLEWORD_MODEL: {os.getenv('DOUBLEWORD_MODEL', 'Qwen/Qwen3-VL-235B-A22B-Instruct-FP8')}")
print(f"  MAX_TOKENS: {os.getenv('MAX_TOKENS', '5000')}")
print()

# Collect files based on arguments
supported_extensions = ['*.pdf', '*.txt', '*.md', '*.docx', '*.pptx', '*.odp']
all_files = []

if args.files:
    # Use specific files provided
    all_files = [str(Path(f).resolve()) for f in args.files]
    print(f"Processing {len(all_files)} specified file(s)\n")
elif args.input_dir:
    # Scan custom directory
    input_dir = Path(args.input_dir)
    if not input_dir.exists():
        print(f"Error: Directory '{args.input_dir}' does not exist")
        exit(1)
    for ext in supported_extensions:
        all_files.extend(glob.glob(str(input_dir / ext)))
    all_files.sort()
    print(f"Found {len(all_files)} files in {args.input_dir}\n")
else:
    # Default: scan data/papers directory
    for ext in supported_extensions:
        all_files.extend(glob.glob(f'data/papers/{ext}'))
    all_files.sort()
    print(f"Found {len(all_files)} files in data/papers/\n")

requests = []
failed_files = []
extraction_stats = {'pypdf': 0, 'pdfplumber': 0, 'txt': 0, 'docx': 0, 'pptx': 0, 'odp': 0}

def extract_text_pypdf(pdf_path):
    """Try pypdf first (faster)."""
    with open(pdf_path, 'rb') as f:
        reader = PdfReader(f)
        text = '\n'.join(page.extract_text() for page in reader.pages)
        return text, len(reader.pages)

def extract_text_pdfplumber(pdf_path):
    """Fallback to pdfplumber (more robust but slower)."""
    with pdfplumber.open(pdf_path) as pdf:
        text = '\n'.join((page.extract_text() or '') for page in pdf.pages)
        return text, len(pdf.pages)

def extract_from_text(file_path):
    """Extract text from .txt or .md files."""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        text = f.read()
        return text, 1

def extract_from_docx(file_path):
    """Extract text from .docx files."""
    doc = Document(file_path)
    paragraphs = [para.text for para in doc.paragraphs]
    text = '\n'.join(paragraphs)
    # Estimate pages (rough: 500 words per page)
    word_count = len(text.split())
    pages = max(1, word_count // 500)
    return text, pages

def extract_from_pptx(file_path):
    """Extract text from .pptx files."""
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    text = '\n'.join(text_runs)
    return text, len(prs.slides)

def extract_from_odp(file_path):
    """Extract text from .odp files."""
    doc = load_odf(file_path)
    text_runs = []
    # Extract all text paragraphs
    for paragraph in doc.getElementsByType(P):
        text_content = ''.join(node.data for node in paragraph.childNodes if hasattr(node, 'data'))
        if text_content.strip():
            text_runs.append(text_content)
    text = '\n'.join(text_runs)
    # Count frames as slide estimate
    frames = doc.getElementsByType(Frame)
    pages = max(1, len(frames))
    return text, pages

for idx, file_path in enumerate(all_files, 1):
    print(f"[{idx}/{len(all_files)}] Processing {file_path}...")

    text = None
    pages = 0
    extraction_method = None
    file_extension = Path(file_path).suffix.lower()

    try:
        # Route to appropriate extraction method based on file type
        if file_extension == '.pdf':
            # Try pypdf first (faster), fallback to pdfplumber
            try:
                text, pages = extract_text_pypdf(file_path)
                extraction_method = 'pypdf'
                extraction_stats['pypdf'] += 1
            except (KeyError, Exception) as e:
                if 'bbox' in str(e) or isinstance(e, KeyError):
                    print(f"  ⚠ pypdf failed ({e}), trying pdfplumber...")
                    text, pages = extract_text_pdfplumber(file_path)
                    extraction_method = 'pdfplumber'
                    extraction_stats['pdfplumber'] += 1
                else:
                    raise

        elif file_extension == '.docx':
            text, pages = extract_from_docx(file_path)
            extraction_method = 'docx'
            extraction_stats['docx'] += 1

        elif file_extension == '.pptx':
            text, pages = extract_from_pptx(file_path)
            extraction_method = 'pptx'
            extraction_stats['pptx'] += 1

        elif file_extension == '.odp':
            text, pages = extract_from_odp(file_path)
            extraction_method = 'odp'
            extraction_stats['odp'] += 1

        elif file_extension in ['.txt', '.md']:
            text, pages = extract_from_text(file_path)
            extraction_method = 'txt'
            extraction_stats['txt'] += 1

        else:
            print(f"  ⚠ Unsupported file type: {file_extension}")
            failed_files.append((file_path, f"unsupported file type: {file_extension}"))
            continue

    except Exception as e:
        print(f"  ✗ Error: {e}")
        failed_files.append((file_path, str(e)))
        continue

    # Skip if no meaningful text extracted
    if not text or len(text.strip()) < 100:
        print(f"  ⚠ Skipped (insufficient text: {len(text)} chars)")
        failed_files.append((file_path, "insufficient text"))
        continue

    print(f"  ✓ Extracted {len(text)} characters from {pages} pages [{extraction_method}]")

    # Create batch request with sanitized custom_id
    # Remove special chars from filename for custom_id (max 64 chars including 'summary-' prefix)
    safe_filename = Path(file_path).stem.replace('%', '_').replace(' ', '_').replace('&', 'and')[:55]

    request = {
        "custom_id": f"summary-{safe_filename}",
        "method": "POST",
        "url": os.getenv('CHAT_COMPLETIONS_ENDPOINT', '/v1/chat/completions'),
        "body": {
            "model": os.getenv('DOUBLEWORD_MODEL', 'Qwen/Qwen3-VL-235B-A22B-Instruct-FP8'),
            "messages": [
                {
                    "role": "user",
                    "content": f"{prompt_template}\n\nDocument text:\n{text}"
                }
            ],
            "max_tokens": int(os.getenv('MAX_TOKENS', '5000'))
        }
    }
    requests.append(request)

# Write JSONL file with timestamp
timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
output_file = f'batch_requests_{timestamp}.jsonl'
with open(output_file, 'w') as f:
    for req in requests:
        f.write(json.dumps(req) + '\n')

print(f"\n{'='*60}")
print(f"✓ Created {output_file} with {len(requests)} requests")
print(f"\nExtraction methods used:")
for method, count in extraction_stats.items():
    if count > 0:
        label = "pdfplumber (fallback)" if method == 'pdfplumber' else method
        print(f"  {label}: {count} files")

if failed_files:
    print(f"\n⚠ Failed to process {len(failed_files)} files:")
    for path, reason in failed_files:
        print(f"  - {Path(path).name}: {reason}")

print(f"\nNext step: python submit_batch.py")
